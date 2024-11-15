from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import Table, TableStyle
import pdfplumber
import markdown2
import os
from pdf2docx import Converter
from PIL import Image as PILImage
import speech_recognition as sr
from pydub import AudioSegment
from pptx import Presentation
import io
import fitz  # PyMuPDF for PDF compression
import gzip
import shutil
import subprocess


def docx_to_pdf(input_file, output_file):
    doc = Document(input_file)
    c = canvas.Canvas(output_file, pagesize=A4)
    width, height = A4
    y_position = height - inch

    c.setFont("Helvetica", 12)

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            c.drawString(1 * inch, y_position, text)
            y_position -= 15

            if y_position < inch:
                c.showPage()
                y_position = height - inch
                c.setFont("Helvetica", 12)

    for table in doc.tables:
        data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        col_count = len(data[0]) if data else 1
        col_width = (width - 2 * inch) / col_count

        pdf_table = Table(data, colWidths=[col_width] * col_count)
        pdf_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ]))

        table_width, table_height = pdf_table.wrap(width - 2 * inch, height)
        if y_position - table_height < inch:
            c.showPage()
            y_position = height - inch

        pdf_table.drawOn(c, 1 * inch, y_position - table_height)
        y_position -= (table_height + 20)

    c.save()
    print(f"Converted {input_file} to {output_file}")


def pdf_to_text(input_file, output_file):
    text = ""
    try:
        with pdfplumber.open(input_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return
    
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"Converted {input_file} to {output_file}")


def markdown_to_html(input_file, output_file):
    with open(input_file, "r", encoding="utf-8") as f:
        md_content = f.read()
    html_content = markdown2.markdown(md_content)
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(html_content)
    print(f"Converted {input_file} to {output_file}")


def txt_to_pdf(input_file, output_file):
    c = canvas.Canvas(output_file, pagesize=A4)
    width, height = A4
    y_position = height - inch

    c.setFont("Helvetica", 12)

    with open(input_file, "r", encoding="utf-8") as f:
        for line in f:
            c.drawString(1 * inch, y_position, line.strip())
            y_position -= 12

            if y_position < inch:
                c.showPage()
                y_position = height - inch
                c.setFont("Helvetica", 12)

    c.save()
    print(f"Converted {input_file} to {output_file}")


def pdf_to_docx(input_file, output_file):
    cv = Converter(input_file)
    cv.convert(output_file)
    cv.close()
    print(f"Converted {input_file} to {output_file}")


def image_to_pdf(image_files, output_file):
    c = canvas.Canvas(output_file, pagesize=A4)
    width, height = A4

    for img_file in image_files:
        img = PILImage.open(img_file)
        
        # Maintain aspect ratio
        img_ratio = min(width / img.width, height / img.height)
        img_width = int(img.width * img_ratio)
        img_height = int(img.height * img_ratio)

        # Center the image on the page
        x_position = (width - img_width) / 2
        y_position = (height - img_height) / 2

        # Draw the image centered on the canvas
        c.drawImage(img_file, x_position, y_position, width=img_width, height=img_height)
        c.showPage()

    c.save()
    print(f"Converted images to {output_file}")


def audio_to_text(input_file, output_file):
    recognizer = sr.Recognizer()
    if input_file.endswith('.mp3'):
        sound = AudioSegment.from_mp3(input_file)
        temp_file = "temp_audio.wav"
        sound.export(temp_file, format="wav")
        input_file = temp_file

    with sr.AudioFile(input_file) as source:
        audio = recognizer.record(source)
        text = recognizer.recognize_google(audio)

    with open(output_file, "w") as f:
        f.write(text)
    print(f"Transcribed audio to {output_file}")

    if input_file == "temp_audio.wav":
        os.remove("temp_audio.wav")

def pptx_to_pdf(input_file, output_file):
    prs = Presentation(input_file)
    c = canvas.Canvas(output_file, pagesize=A4)
    width, height = A4

    for slide in prs.slides:
        y_position = height - 40  

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                  
                    lines = text.splitlines()
                    for line in lines:
                        c.drawString(40, y_position, line)
                        y_position -= 12
                        if y_position < 40:
                            c.showPage()
                            y_position = height - 40

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image_stream = io.BytesIO(shape.image.blob)
                img = PILImage.open(image_stream)
                img_ratio = min(width / img.width, (y_position - 40) / img.height)
                img_width, img_height = int(img.width * img_ratio), int(img.height * img_ratio)

            
                c.drawImage(image_stream, (width - img_width) / 2, y_position - img_height, 
                            width=img_width, height=img_height)
                y_position -= (img_height + 20)  

        c.showPage()  

    c.save()
    print(f"Converted {input_file} to {output_file}")

def compress_pdf(input_file, output_file, method="pymupdf", quality=20, zip_output=False):
    try:
        if method == "pymupdf":
            pdf_document = fitz.open(input_file)
            pdf_document.save(output_file, garbage=4, deflate=True, clean=True, pretty=True)
            pdf_document.close()
            print(f"Compressed PDF saved as {output_file} using PyMuPDF.")

        elif method == "ghostscript":
            command = [
                "gs",
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/screen",
                "-dNOPAUSE",
                "-dBATCH",
                "-sOutputFile=" + output_file,
                input_file,
            ]
            subprocess.run(command, check=True)
            print(f"Compressed PDF saved as {output_file} using Ghostscript.")

        else:
            raise ValueError("Invalid method. Use 'pymupdf' or 'ghostscript'.")
        if zip_output:
            zip_file_path = output_file + ".zip"
            with shutil.ZipFile(zip_file_path, 'w') as zipf:
                zipf.write(output_file, os.path.basename(output_file))
            print(f"Compressed PDF saved as {zip_file_path} in a .zip archive.")

    except Exception as e:
        print(f"Error compressing PDF: {e}")



def compress_text_file(input_file, output_file, method="gzip"):

    try:
        if method == "gzip":
            
            if not output_file.endswith(".gz"):
                output_file += ".gz"
            # Compress
            with open(input_file, 'rb') as f_in:
                with gzip.open(output_file, 'wb') as f_out:
                    shutil.copyfileobj(f_in, f_out)
            print(f"Compressed text file saved as {output_file} using Gzip.")
        
        elif method == "zip":
       
            if not output_file.endswith(".zip"):
                output_file += ".zip"
           
            shutil.make_archive(output_file.replace(".zip", ""), 'zip', os.path.dirname(input_file), os.path.basename(input_file))
            print(f"Compressed text file saved as {output_file} using Zip.")
        
        else:
            print("Error: Invalid method. Use 'gzip' or 'zip'.")
    
    except Exception as e:
        print(f"Error compressing file: {e}")